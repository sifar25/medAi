from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent.parent
OUTPUT = ROOT / "presentations" / "MedAgent_Supervisor_Deck.pptx"
SCREENSHOT_DIR = ROOT / "presentations" / "screenshots"
STATUS_OUTPUT = SCREENSHOT_DIR / "status_output.txt"

NAVY = RGBColor(8, 37, 53)
TEAL = RGBColor(0, 95, 115)
GOLD = RGBColor(238, 155, 0)
SAND = RGBColor(245, 239, 228)
INK = RGBColor(31, 42, 55)
MUTED = RGBColor(94, 109, 126)
WHITE = RGBColor(255, 255, 255)
PANEL = RGBColor(251, 248, 242)
SOFT = RGBColor(235, 245, 246)
LINE = RGBColor(232, 223, 210)


SLIDES = [
    {
        "title": "MedAgent: Human-Centred AI Support for Medication Adherence in Chronic Illness",
        "subtitle": "MSc prototype deck for supervisor meetings, viva preparation, and project review",
        "points": [
            "Early identification of medication non-adherence risk in chronic illness follow-up",
            "Local, reproducible prototype built with synthetic data only",
            "Clinician-supporting workflow rather than autonomous decision-making",
        ],
        "notes": "Open with the problem clearly: clinicians need earlier, more structured visibility of adherence risk without losing human judgment.",
        "tag": "Slide 1 | Project framing",
        "image": "dashboard.png",
    },
    {
        "title": "Why This Problem Matters",
        "points": [
            "Missed medication can worsen long-term outcomes and destabilise chronic care",
            "Non-adherence is shaped by routine, side effects, behaviour, and follow-up barriers",
            "Research goal: combine prediction with practical, human-readable action support",
        ],
        "notes": "Explain that the project is not only about predicting a label. The goal is to make the output usable in practice by translating risk into readable summaries and realistic follow-up actions.",
        "tag": "Slide 2 | Motivation",
        "image": "login.png",
        "image_label": "Professional local login boundary for demo use",
    },
    {
        "title": "Prototype Architecture Overview",
        "points": [
            "Flask web interface for login, dashboard, assessment, and model maintenance",
            "Random Forest model for medication adherence risk prediction",
            "Agent layer for explanation and support-plan generation",
            "CSV-backed synthetic dataset for simple, local reproducible execution",
        ],
        "notes": "Walk from input to output: user interaction, model prediction, then agent-generated explanation and support guidance.",
        "tag": "Slide 3 | Architecture",
        "diagram": True,
    },
    {
        "title": "Clinician-Facing Dashboard",
        "points": [
            "Summary cards provide a quick view of case distribution across risk levels",
            "Search, filtering, and pagination support structured review of the patient list",
            "The dashboard acts as the triage entry point into the system",
        ],
        "notes": "Use this slide to explain the move from a rough prototype table to a more realistic case-review workflow.",
        "tag": "Slide 4 | Dashboard",
        "image": "dashboard.png",
        "image_label": "Population-level overview with case summary and triage entry point",
    },
    {
        "title": "Usability Improvements for Review Workflows",
        "points": [
            "Filtered dashboard states show targeted exploration by risk and condition",
            "Pagination improves readability and keeps the interface presentation-ready",
            "UI polish supports clearer supervisor review and MSc demonstration quality",
        ],
        "notes": "Justify the polish work here: these changes improve both usability and the credibility of the project demonstration.",
        "tag": "Slide 5 | UX polish",
        "image": "dashboard_filtered.png",
        "image_label": "Filtered dashboard showing targeted case review",
    },
    {
        "title": "From Risk Prediction to Human-Readable Guidance",
        "points": [
            "Risk level and confidence are shown clearly at patient level",
            "Explanatory factors summarise why a patient may be concerning",
            "Support recommendations translate analytics into follow-up actions",
            "High-risk cases explicitly retain clinician oversight",
        ],
        "notes": "Emphasise that the system does not stop at prediction. It also helps interpret the result in a way that is easier to discuss and act on.",
        "tag": "Slide 6 | Patient result",
        "image": "patient_result.png",
        "image_label": "Patient assessment screen with risk panel and support plan",
    },
    {
        "title": "Running a New Assessment Live",
        "points": [
            "New patient scenarios can be entered directly through the web form",
            "Optional persistence allows demo records to be kept in the working dataset",
            "Duplicate protection prevents repeated accidental insertion of the same case",
        ],
        "notes": "Describe how a live demo works: enter a case, generate a result, optionally save it, then show duplicate protection on re-submit.",
        "tag": "Slide 7 | Live workflow",
        "image": "new_assessment_form.png",
        "secondary_image": "new_assessment_result.png",
        "image_label": "Form entry and resulting assessment output",
    },
    {
        "title": "Reproducibility, Idempotency, and Local Execution",
        "points": [
            "Health-checked startup and controlled shutdown scripts support reliable demonstrations",
            "Fresh-start and reseed controls allow deterministic dataset restoration",
            "Default restart behaviour preserves user-added records unless reset is explicit",
            "Windows, macOS, and Ubuntu execution guidance is documented",
        ],
        "notes": "Use this slide to defend engineering quality. The prototype is runnable, testable, and repeatable rather than a single-use demo.",
        "tag": "Slide 8 | Operations",
        "image": "train_page.png",
        "image_label": "Model maintenance remains local and reproducible",
    },
    {
        "title": "What Was Tested",
        "points": [
            "Shell smoke validation covers login, dashboard, assessment, save, and dedupe flow",
            "Playwright browser tests cover lifecycle behaviour and UI interaction",
            "The current build passed both shell and browser validation",
        ],
        "notes": "Keep this factual: the project was validated with executable checks, not only manual clicking.",
        "tag": "Slide 9 | Validation",
        "status_text": True,
    },
    {
        "title": "Contribution and Next Development Steps",
        "points": [
            "Contribution: predictive triage combined with human-readable support recommendations",
            "Current limitation: synthetic data and local prototype scope only",
            "Next steps: richer analytics, report outputs, and broader clinical validation",
        ],
        "notes": "Close with a balanced evaluation: be clear about the contribution, but equally clear about current prototype boundaries.",
        "tag": "Slide 10 | Closing",
        "image": "dashboard.png",
        "image_label": "Final prototype state ready for supervisor review",
    },
]


def set_background(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_header_band(slide, text):
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.4))
    band.fill.solid()
    band.fill.fore_color.rgb = GOLD
    band.line.fill.background()
    frame = band.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = NAVY


def add_footer_tag(slide, text):
    box = slide.shapes.add_textbox(Inches(10.35), Inches(7.0), Inches(2.55), Inches(0.24))
    frame = box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.RIGHT
    p.font.size = Pt(10)
    p.font.color.rgb = MUTED


def add_text_box(slide, left, top, width, height, text, size=16, bold=False, color=INK):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return box


def add_bullet_card(slide, spec):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.28), Inches(4.15), Inches(5.95))
    card.fill.solid()
    card.fill.fore_color.rgb = PANEL
    card.line.color.rgb = LINE
    frame = card.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.2)
    frame.margin_right = Inches(0.16)
    frame.margin_top = Inches(0.18)
    title = frame.paragraphs[0]
    title.text = "Key message"
    title.font.size = Pt(14)
    title.font.bold = True
    title.font.color.rgb = TEAL
    for bullet in spec["points"]:
        p = frame.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(16)
        p.font.color.rgb = INK
        p.space_after = Pt(8)


def fit_picture(slide, image_path, left, top, width, height):
    with Image.open(image_path) as image:
        img_w, img_h = image.size
    box_ratio = float(width) / float(height)
    img_ratio = img_w / img_h

    if img_ratio > box_ratio:
        scaled_width = height * img_ratio
        scaled_height = height
    else:
        scaled_width = width
        scaled_height = width / img_ratio

    pic_left = left + (width - scaled_width) / 2
    pic_top = top + (height - scaled_height) / 2
    slide.shapes.add_picture(str(image_path), pic_left, pic_top, width=scaled_width, height=scaled_height)


def add_image_card(slide, image_name, left, top, width, height, label=None):
    image_path = SCREENSHOT_DIR / image_name
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = LINE
    if image_path.exists():
        fit_picture(slide, image_path, left + Inches(0.08), top + Inches(0.08), width - Inches(0.16), height - Inches(0.44))
    else:
        frame = card.text_frame
        frame.clear()
        p = frame.paragraphs[0]
        p.text = f"Missing image: {image_name}"
        p.font.size = Pt(14)
        p.font.color.rgb = MUTED
    if label:
        label_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left + Inches(0.2), top + height - Inches(0.42), width - Inches(0.4), Inches(0.26))
        label_box.fill.solid()
        label_box.fill.fore_color.rgb = NAVY
        label_box.line.fill.background()
        frame = label_box.text_frame
        frame.clear()
        p = frame.paragraphs[0]
        p.text = label
        p.font.size = Pt(10)
        p.font.color.rgb = WHITE


def add_notes_card(slide, spec, left=Inches(4.95), top=Inches(5.55), width=Inches(3.3), height=Inches(1.4)):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = SOFT
    card.line.color.rgb = RGBColor(205, 227, 229)
    frame = card.text_frame
    frame.clear()
    frame.word_wrap = True
    title = frame.paragraphs[0]
    title.text = "Speaker note"
    title.font.size = Pt(12)
    title.font.bold = True
    title.font.color.rgb = TEAL
    p = frame.add_paragraph()
    p.text = spec["notes"]
    p.font.size = Pt(11)
    p.font.color.rgb = INK


def add_architecture_diagram(slide):
    labels = [
        ("Data layer", "Synthetic patient records\npatients_seed.csv -> patients.csv", Inches(5.1)),
        ("Prediction layer", "Random Forest model\nprediction_model.py", Inches(7.15)),
        ("Agent layer", "Explanation + support plan\nagent.py", Inches(9.2)),
    ]
    for title, body, left in labels:
        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, Inches(2.15), Inches(1.7), Inches(2.0))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = LINE
        frame = box.text_frame
        frame.clear()
        frame.word_wrap = True
        p = frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = NAVY
        p = frame.add_paragraph()
        p.text = body
        p.font.size = Pt(12)
        p.font.color.rgb = INK

    input_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(5.0), Inches(4.75), Inches(2.35), Inches(0.9))
    input_box.fill.solid()
    input_box.fill.fore_color.rgb = GOLD
    input_box.line.fill.background()
    frame = input_box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = "Clinician inputs and patient context"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = NAVY

    output_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.15), Inches(4.75), Inches(2.75), Inches(0.9))
    output_box.fill.solid()
    output_box.fill.fore_color.rgb = TEAL
    output_box.line.fill.background()
    frame = output_box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = "Risk label, explanation, and support plan"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE


def add_status_card(slide):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(4.95), Inches(1.35), Inches(7.55), Inches(4.2))
    card.fill.solid()
    card.fill.fore_color.rgb = NAVY
    card.line.fill.background()
    frame = card.text_frame
    frame.clear()
    frame.word_wrap = True
    title = frame.paragraphs[0]
    title.text = "Verification evidence"
    title.font.size = Pt(15)
    title.font.bold = True
    title.font.color.rgb = WHITE
    status_text = "Status output unavailable"
    if STATUS_OUTPUT.exists():
        status_text = STATUS_OUTPUT.read_text(encoding="utf-8").strip()
    body = frame.add_paragraph()
    body.text = status_text
    body.font.size = Pt(13)
    body.font.name = "Courier New"
    body.font.color.rgb = RGBColor(224, 241, 242)

    badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(9.65), Inches(5.85), Inches(2.55), Inches(0.52))
    badge.fill.solid()
    badge.fill.fore_color.rgb = GOLD
    badge.line.fill.background()
    frame = badge.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = "Smoke + browser QA passed"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = NAVY


def add_title_slide(prs, spec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, NAVY)
    image_path = SCREENSHOT_DIR / spec["image"]
    if image_path.exists():
        fit_picture(slide, image_path, Inches(0), Inches(0), Inches(13.333), Inches(7.5))

    overlay = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = NAVY
    overlay.fill.transparency = 0.36
    overlay.line.fill.background()

    banner = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(0.65), Inches(12.0), Inches(5.9))
    banner.fill.solid()
    banner.fill.fore_color.rgb = RGBColor(9, 51, 67)
    banner.fill.transparency = 0.1
    banner.line.fill.background()

    add_text_box(slide, Inches(1.0), Inches(1.05), Inches(10.8), Inches(1.35), spec["title"], size=25, bold=True, color=WHITE)
    add_text_box(slide, Inches(1.0), Inches(2.45), Inches(10.0), Inches(0.5), spec["subtitle"], size=16, color=RGBColor(224, 241, 242))

    callout = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(4.55), Inches(5.45), Inches(1.45))
    callout.fill.solid()
    callout.fill.fore_color.rgb = WHITE
    callout.line.fill.background()
    frame = callout.text_frame
    frame.clear()
    frame.word_wrap = True
    for index, bullet in enumerate(spec["points"]):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.text = bullet
        p.font.size = Pt(15)
        p.font.color.rgb = INK

    note = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.75), Inches(4.55), Inches(4.7), Inches(1.45))
    note.fill.solid()
    note.fill.fore_color.rgb = SAND
    note.line.fill.background()
    frame = note.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = "Presentation focus"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p = frame.add_paragraph()
    p.text = spec["notes"]
    p.font.size = Pt(12)
    p.font.color.rgb = INK

    add_footer_tag(slide, spec["tag"])


def add_standard_slide(prs, spec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, WHITE)
    add_header_band(slide, "MedAgent supervisor presentation")
    add_text_box(slide, Inches(0.65), Inches(0.7), Inches(7.5), Inches(0.5), spec["title"], size=24, bold=True, color=NAVY)
    add_bullet_card(slide, spec)

    if spec.get("diagram"):
        add_architecture_diagram(slide)
        add_notes_card(slide, spec)
    elif spec.get("status_text"):
        add_status_card(slide)
        add_notes_card(slide, spec, left=Inches(4.95), top=Inches(5.85), width=Inches(4.25), height=Inches(1.0))
    elif spec.get("secondary_image"):
        add_image_card(slide, spec["image"], Inches(4.95), Inches(1.45), Inches(3.55), Inches(2.35), "Scenario entry")
        add_image_card(slide, spec["secondary_image"], Inches(8.75), Inches(1.45), Inches(3.55), Inches(2.35), "Result with persistence")
        add_notes_card(slide, spec, left=Inches(4.95), top=Inches(4.15), width=Inches(7.35), height=Inches(1.45))
    else:
        add_image_card(slide, spec["image"], Inches(4.95), Inches(1.45), Inches(7.35), Inches(4.15), spec.get("image_label"))
        add_notes_card(slide, spec)

    add_footer_tag(slide, spec["tag"])


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs, SLIDES[0])
    for spec in SLIDES[1:]:
        add_standard_slide(prs, spec)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(f"Created {OUTPUT}")


if __name__ == "__main__":
    build_deck()